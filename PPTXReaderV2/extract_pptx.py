'''
Created on Jan 16, 2019

@author: Bowonwit
'''

from __future__ import print_function, unicode_literals
from pptx import Presentation
from PyInquirer import prompt, print_json
import os
import time
from subF_titleO import *

def extract_file(path_file, file_name):
    print (path_file+file_name)
    prs = Presentation(path_file+file_name)
    ido_obs = [titleO(-1,0,False,file_name)]
    temp_title_obs = find_title(prs,ido_obs[0])
    ido_obs = []
    
    for _tto in temp_title_obs:
        ido_obs_buffer = []
        for i,slide in enumerate(prs.slides):
            if i+1 in _tto['pages']:
                shapes = resort_shapes(slide.shapes)
                for shape in shapes:
                    if shape.has_text_frame and not(replace_for_compare(slide.shapes.title.text) == replace_for_compare(shape.text)):
                        for paragraph in shape.text_frame.paragraphs:
                            if replace_2019(paragraph.text) != "" :#and replace_2019(paragraph.text) not in [iob[3] for iob in ido_obs_buffer]:
                                ido_obs_buffer.append([i+1,paragraph.level,tap_counter(replace_2019(paragraph.text)),replace_2019(paragraph.text)])
                                
        ido_obs.append(_tto['object'])
        temp_last_obj = [0,_tto['object']]
        
        '''c
        for s in ido_obs_buffer:
            print([s[0]],end=",___")
            print([s[1]],end=",___")
            print([s[2]],end=",___")
            print([s[3]])
        '''
        
        for i,_iob in enumerate(ido_obs_buffer):
            
                
            if i == 0:
                _ctr = titleO(temp_last_obj[1].get_io_type()+1,_iob[0],temp_last_obj[1],_iob[3])
                ido_obs.append(_ctr)
                temp_last_obj = [_iob[1],_ctr]
            else:
                if temp_last_obj[0] < _iob[1]:
                    _ctr = titleO(temp_last_obj[1].get_io_type()+1,_iob[0],temp_last_obj[1],_iob[3])
                    ido_obs.append(_ctr)
                    temp_last_obj = [_iob[1],_ctr]

                elif temp_last_obj[0] == _iob[1]:
                    _ctr = titleO(temp_last_obj[1].get_io_parent().get_io_type()+1,_iob[0],temp_last_obj[1].get_io_parent(),_iob[3])
                    ido_obs.append(_ctr)
                    temp_last_obj = [_iob[1],_ctr]
                else:
                    check = last_priority_obj(ido_obs, ido_obs_buffer, _iob[1], i)
                    if check == False:
                        _ctr = titleO(ido_obs[0].get_io_type()+1,_iob[0],ido_obs[0],_iob[3])
                        ido_obs.append(_ctr)
                        temp_last_obj = [_iob[1],_ctr]
                    else:
                        _ctr = titleO(check.get_io_type()+1,_iob[0],check,_iob[3])
                        ido_obs.append(_ctr)
                        temp_last_obj = [_iob[1],_ctr]
    
    reduce_redundance(ido_obs)
    
    '''
    for ds in ido_obs:
        print(ds.get_io_page(),end=",___")
        print(ds.get_io_type(),end=",___")
        if ds.get_io_parent() != False :print(ds.get_io_parent().get_io_type(),end=",___")
        print(ds.get_io_data())
    '''
    
    return ido_obs

def tour_in_pptx(ido_obs, parent, _pri):
    while True:
        name = 'title'
        if _pri != 0: name = 'body'
        list_que, list_obs = create_que_T_list(ido_obs,find_by_prior(ido_obs, parent, _pri))
        list_que.append("Back...")
        que_title = [
                        {
                            'type': 'list',
                            'name': 'title',
                            'message': 'List of ' + name + ' in PPTX file...',
                            'choices': list_que,
                        },
                    ]
        aws_title = prompt(que_title)
        if aws_title['title'] == que_title[0]['choices'][-1]: break
        else:
            if find_parent(ido_obs, list_obs[find_index_que(aws_title['title'], que_title[0]['choices'])]) == False:
                tour_in_pptx(ido_obs, parent, _pri)
            else:
                tour_in_pptx(ido_obs, find_parent(ido_obs, list_obs[find_index_que(aws_title['title'], que_title[0]['choices'])]), find_parent(ido_obs, list_obs[find_index_que(aws_title['title'], que_title[0]['choices'])]).get_io_type()+1)

def handle_extract(path_file, file_name):
    ido_obs = extract_file(path_file, file_name)
    tour_in_pptx(ido_obs,ido_obs[0], 0)


            